from flask import Blueprint, render_template, request, redirect, url_for, send_file, flash
import pptx
from pptx.util import Inches, Pt
from io import BytesIO

pptx_blueprint = Blueprint('pptx_blueprint', __name__)

# Fungsi untuk memecah teks menjadi beberapa halaman sesuai batasan baris
def split_text_to_slides(text, max_lines_per_slide=12, max_words_per_line=10):
    words = text.split()  # Memisahkan teks menjadi kata
    slides = []
    current_slide = []
    current_line = []

    for word in words:
        # Tambahkan kata ke baris saat ini
        current_line.append(word)

        # Jika panjang baris lebih dari batasan, tambahkan baris ke slide dan reset baris saat ini
        if len(current_line) > max_words_per_line:
            current_slide.append(' '.join(current_line[:-1]))  # Tambahkan semua kecuali kata terakhir
            current_line = [current_line[-1]]  # Mulai baris baru dengan kata terakhir

        # Jika jumlah baris dalam slide mencapai batas, tambahkan slide dan reset slide saat ini
        if len(current_slide) >= max_lines_per_slide:
            slides.append('\n'.join(current_slide))  # Tambahkan slide yang sudah lengkap
            current_slide = []
    
    # Tambahkan baris terakhir yang tersisa ke slide
    if current_line:
        current_slide.append(' '.join(current_line))
    if current_slide:
        slides.append('\n'.join(current_slide))

    return slides

# Fungsi untuk membuat PPTX berdasarkan teks dan footer
def create_pptx(title, text_slides, footer_text):
    prs = pptx.Presentation()

    for slide_text in text_slides:
        slide_layout = prs.slide_layouts[1]  # Layout kedua: Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Menambahkan judul
        title_box = slide.shapes.title
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(20)  # Ukuran font untuk judul

        # Menambahkan isi teks
        content_box = slide.shapes.placeholders[1]
        content_box.text = slide_text
        
        # Mengatur ukuran font untuk semua isi teks
        for paragraph in content_box.text_frame.paragraphs:
            paragraph.font.size = Pt(18)  # Ukuran font untuk isi teks

        # Menambahkan footer
        left = Inches(0.5)
        top = prs.slide_height - Inches(1)
        width = prs.slide_width - Inches(1)
        height = Inches(0.5)

        footer_box = slide.shapes.add_textbox(left, top, width, height)
        footer_frame = footer_box.text_frame
        footer_paragraph = footer_frame.add_paragraph()
        footer_paragraph.text = footer_text
        footer_paragraph.font.size = Pt(8)  # Ukuran font untuk footer

    # Simpan PPTX ke memory buffer (BytesIO)
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    return ppt_buffer

# Halaman pilihan format PPTX
@pptx_blueprint.route('/pptx_choice', methods=['GET', 'POST'])
def pptx_choice():
    if request.method == 'POST':
        title = request.form.get('title')  # Mendapatkan judul dari form
        footer = request.form.get('footer')  # Mendapatkan footer dari form
        text = request.form.get('text')  # Mendapatkan teks dari form

        # Validasi input
        if not title or not text:
            flash('Judul dan teks tidak boleh kosong!', 'error')
            return redirect(url_for('pptx_blueprint.pptx_choice'))

        # Pisahkan teks menjadi beberapa halaman (slide)
        text_slides = split_text_to_slides(text)

        # Buat file PPTX berdasarkan teks
        pptx_file = create_pptx(title, text_slides, footer)

        # Kirim file PPTX sebagai unduhan
        return send_file(pptx_file, download_name="translated_presentation.pptx", as_attachment=True)

    # Mengambil teks terjemahan dari request sebelumnya jika ada
    translated_text = request.args.get('translated_text', '')  # Ambil teks terjemahan jika ada
    return render_template('pptx_choice.html', translated_text=translated_text)
