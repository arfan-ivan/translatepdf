import re
from flask import Flask, render_template, request, redirect, url_for, send_file
from werkzeug.utils import secure_filename
import os
import requests  # Pastikan pustaka ini diimpor
import PyPDF2
import docx
import odf.opendocument
import odf.text
import pptx
from fpdf import FPDF
from pptx_handler import pptx_blueprint  # Sesuaikan nama file jika perlu
app = Flask(__name__)

app.register_blueprint(pptx_blueprint, url_prefix='/pptx')  # Pastikan url_prefix sesuai


app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Supported languages for translation with user-friendly names
SUPPORTED_LANGUAGES = {
    'al': 'Albanian',
    'ar': 'Arabic',
    'az': 'Azerbaijani',
    'eu': 'Basque',
    'bn': 'Bengali',
    'bg': 'Bulgarian',
    'ca': 'Catalan',
    'zh': 'Chinese',
    'zh-tw': 'Chinese (Traditional)',
    'cs': 'Czech',
    'da': 'Danish',
    'nl': 'Dutch',
    'en': 'English',
    'eo': 'Esperanto',
    'et': 'Estonian',
    'fi': 'Finnish',
    'fr': 'French',
    'gl': 'Galician',
    'de': 'German',
    'el': 'Greek',
    'he': 'Hebrew',
    'hi': 'Hindi',
    'hu': 'Hungarian',
    'id': 'Indonesian',
    'ga': 'Irish',
    'it': 'Italian',
    'ja': 'Japanese',
    'ko': 'Korean',
    'lv': 'Latvian',
    'lt': 'Lithuanian',
    'ml': 'Malay',
    'no': 'Norwegian',
    'fa': 'Persian',
    'pl': 'Polish',
    'pt': 'Portuguese',
    'ro': 'Romanian',
    'ru': 'Russian',
    'sk': 'Slovak',
    'sl': 'Slovenian',
    'es': 'Spanish',
    'sv': 'Swedish',
    'tl': 'Tagalog',
    'th': 'Thai',
    'tr': 'Turkish',
    'uk': 'Ukrainian',
    'ur': 'Urdu'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'pdf', 'txt', 'odt', 'docx', 'pptx'}

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_odt(file_path):
    doc = odf.opendocument.load(file_path)
    text = ""
    for element in doc.getElementsByType(odf.text.P):
        text += element.firstChild.nodeValue + "\n"
    return text

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path):
    ext = file_path.rsplit('.', 1)[1].lower()
    if ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif ext == 'docx':
        return extract_text_from_docx(file_path)
    elif ext == 'odt':
        return extract_text_from_odt(file_path)
    elif ext == 'pptx':
        return extract_text_from_pptx(file_path)
    elif ext == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    return ""

def save_text_as_pdf(text, filename):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for line in text.splitlines():
        pdf.cell(200, 10, txt=line.encode('latin-1', 'replace').decode('latin-1'), ln=True)
    pdf.output(filename)

def save_text_as_docx(text, filename):
    doc = docx.Document()
    for line in text.splitlines():
        doc.add_paragraph(line)
    doc.save(filename)

def save_text_as_txt(text, filename):
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(text)

@app.route('/')
def upload_file():
    return render_template('upload.html', languages=SUPPORTED_LANGUAGES)

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return redirect(request.url)

    file = request.files['file']
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        text = extract_text(file_path)
        # Pastikan SUPPORTED_LANGUAGES diteruskan ke template
        return render_template('translate.html', text=text, languages=SUPPORTED_LANGUAGES)

    return redirect(url_for('upload_file'))

@app.route('/translate', methods=['POST'])
def translate():
    text = request.form['text']
    target_lang = request.form['target_lang']

    # Menghilangkan simbol yang mengganggu
    clean_text = re.sub(r'[✧▪]', '', text)  # Hapus simbol yang tidak diinginkan

    # Call the LibreTranslate API
    response = requests.post('http://172.16.1.5:5016/translate', json={
        'q': clean_text,
        'source': 'en',
        'target': target_lang,
        'format': 'text'
    })

    translated_text = response.json().get('translatedText', 'Translation failed.')

    return render_template('preview.html', translated_text=translated_text)

@app.route('/save', methods=['POST'])
def save():
    translated_text = request.form['translated_text']
    format_option = request.form['format_option']
    structured_format = request.form.get('structured_format') == 'yes'  # Cek jika opsi dicentang
    
    if structured_format:
        # Logika untuk menyimpan dengan struktur paragraf
        # Misalnya, memisahkan teks berdasarkan paragraf
        text_to_save = translated_text  # Misalkan kita menggunakan teks yang sudah diproses
    else:
        # Logika untuk menyimpan tanpa struktur paragraf
        text_to_save = translated_text.replace("\n", " ")  # Menghapus newline jika tidak ingin struktur
    
    if format_option == 'pdf':
        filename = 'translated_text.pdf'
        save_text_as_pdf(text_to_save, os.path.join(app.config['UPLOAD_FOLDER'], filename))
    elif format_option == 'docx':
        filename = 'translated_text.docx'
        save_text_as_docx(text_to_save, os.path.join(app.config['UPLOAD_FOLDER'], filename))
    elif format_option == 'txt':
        filename = 'translated_text.txt'
        save_text_as_txt(text_to_save, os.path.join(app.config['UPLOAD_FOLDER'], filename))
    else:
        return "Format tidak dikenal.", 400

    return send_file(os.path.join(app.config['UPLOAD_FOLDER'], filename), as_attachment=True)


@app.route('/ppt_choice', methods=['POST'])
def ppt_choice():
    translated_text = request.form['translated_text']
    return render_template('ppt_choice.html', translated_text=translated_text)


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5017)
